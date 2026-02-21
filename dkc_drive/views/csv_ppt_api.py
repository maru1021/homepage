"""
PowerPoint・Outlook MSG 読み込みAPI
"""
import os
import base64
import logging
import mimetypes
from django.views import View
from django.http import JsonResponse, HttpResponse

from .utils import resolve_filepath, error_response

logger = logging.getLogger(__name__)

# Outlookが危険な拡張子をリネームするマッピング（.exe→.ex_, .zip→.zi_ 等）
OUTLOOK_RENAMED_EXT = {
    '.zi_': '.zip', '.ex_': '.exe', '.ba_': '.bat', '.cm_': '.cmd',
    '.js_': '.js', '.vb_': '.vbs', '.ws_': '.wsf', '.sc_': '.scr',
}

# マジックバイトによるファイル形式判定
MAGIC_BYTES = [
    (b'PK\x03\x04', 'application/zip'),
    (b'\x89PNG', 'image/png'),
    (b'\xff\xd8\xff', 'image/jpeg'),
    (b'GIF8', 'image/gif'),
    (b'%PDF', 'application/pdf'),
    (b'BM', 'image/bmp'),
]


def _guess_content_type(filename, data=None):
    """拡張子 + マジックバイトでMIMEタイプを推定。(content_type, 復元済みfilename)を返す"""
    _, ext = os.path.splitext(filename)
    restored = OUTLOOK_RENAMED_EXT.get(ext.lower())
    if restored:
        filename = filename[:len(filename) - len(ext)] + restored

    ct = mimetypes.guess_type(filename)[0]
    if ct and ct != 'application/octet-stream':
        return ct, filename

    if data and len(data) >= 4:
        for magic, mime in MAGIC_BYTES:
            if data[:len(magic)] == magic:
                return mime, filename

    return 'application/octet-stream', filename


class PowerPointDataAPI(View):
    """PowerPoint読み込みAPI"""

    IMAGE_FORMAT_MAP = {
        'image/png': 'png', 'image/jpeg': 'jpeg', 'image/gif': 'gif',
        'image/bmp': 'bmp', 'image/tiff': 'tiff',
    }

    def get(self, request, *args, **kwargs):
        filepath, _, err = resolve_filepath(request.GET, check_sqlite=False)
        if err:
            return err

        try:
            from pptx import Presentation

            prs = Presentation(filepath)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            slides_data = [
                {'shapes': [s for s in (self._extract_shape(shape, slide_width, slide_height) for shape in slide.shapes) if s]}
                for slide in prs.slides
            ]

            return JsonResponse({'status': 'success', 'slides': slides_data})

        except ImportError:
            return error_response('python-pptxがインストールされていません。pip install python-pptx を実行してください。', status=500)
        except Exception as e:
            return error_response(str(e), status=500)

    def _extract_shape(self, shape, slide_width, slide_height):
        """シェイプから情報を抽出"""
        x = (shape.left / slide_width) * 100 if slide_width else 0
        y = (shape.top / slide_height) * 100 if slide_height else 0
        width = (shape.width / slide_width) * 100 if slide_width else 0
        height = (shape.height / slide_height) * 100 if slide_height else 0
        pos = {'x': round(x, 2), 'y': round(y, 2), 'width': round(width, 2), 'height': round(height, 2)}

        if hasattr(shape, 'text') and shape.text.strip():
            font_size = self._get_font_size(shape)
            return {'type': 'text', 'text': shape.text, **pos, 'fontSize': round(font_size)}

        if hasattr(shape, 'image'):
            try:
                image = shape.image
                img_format = self.IMAGE_FORMAT_MAP.get(image.content_type, 'png')
                return {'type': 'image', 'data': base64.b64encode(image.blob).decode('utf-8'), 'format': img_format, **pos}
            except Exception:
                logger.debug("PowerPoint画像抽出スキップ: %s", getattr(shape, 'name', 'unknown'))

        return None

    def _get_font_size(self, shape):
        """シェイプから最初のフォントサイズを取得（デフォルト14pt）"""
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        return run.font.size.pt
        return 14


class MsgDataAPI(View):
    """Outlook MSG読み込みAPI"""

    def get(self, request, *args, **kwargs):
        filepath, _, err = resolve_filepath(request.GET, check_sqlite=False)
        if err:
            return err

        attachment_index = request.GET.get('attachment')

        try:
            import extract_msg

            with extract_msg.openMsg(filepath) as msg:
                if attachment_index is not None:
                    idx = int(attachment_index)
                    if request.GET.get('preview'):
                        return self._preview_attachment(msg, idx)
                    return self._download_attachment(msg, idx)

                return self._get_mail_info(msg)

        except ImportError:
            return error_response('extract-msgがインストールされていません。pip install extract-msg を実行してください。', status=500)
        except Exception as e:
            logger.exception("MSG読み込みエラー: %s", request.GET.get('path'))
            return error_response(str(e), status=500)

    def _get_mail_info(self, msg):
        """メール情報をJSON形式で返す"""
        attachments = []
        for i, att in enumerate(msg.attachments):
            filename = att.longFilename or att.shortFilename or f'attachment_{i}'
            data = att.data or b''
            content_type, filename = _guess_content_type(filename, data)
            attachments.append({'name': filename, 'size': len(data), 'index': i, 'contentType': content_type})

        return JsonResponse({
            'status': 'success',
            'subject': msg.subject or '',
            'sender': msg.sender or '',
            'to': msg.to or '',
            'cc': msg.cc or '',
            'date': msg.date or '',
            'body': msg.body or '',
            'attachments': attachments,
        })

    def _get_attachment_data(self, msg, index):
        """添付ファイルのデータを取得。失敗時は(None, None, None)"""
        if index < 0 or index >= len(msg.attachments):
            return None, None, None
        att = msg.attachments[index]
        filename = att.longFilename or att.shortFilename or f'attachment_{index}'
        data = att.data or b''
        content_type, filename = _guess_content_type(filename, data)
        return filename, data, content_type

    def _download_attachment(self, msg, index):
        """添付ファイルをダウンロード"""
        filename, data, content_type = self._get_attachment_data(msg, index)
        if filename is None:
            return error_response('添付ファイルが見つかりません', status=404)

        response = HttpResponse(data, content_type=content_type)
        response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response

    def _preview_attachment(self, msg, index):
        """添付ファイルをbase64で返す（モーダルプレビュー用）"""
        filename, data, content_type = self._get_attachment_data(msg, index)
        if filename is None:
            return error_response('添付ファイルが見つかりません', status=404)

        return JsonResponse({
            'status': 'success',
            'name': filename,
            'contentType': content_type,
            'data': base64.b64encode(data).decode('utf-8'),
        })
